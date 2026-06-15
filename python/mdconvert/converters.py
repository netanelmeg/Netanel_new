"""Format-specific handlers for mdconvert.

Each handler is registered against one or more file extensions via the
``@register`` decorator from :mod:`mdconvert.core`. Stdlib-only formats (text,
CSV/TSV, JSON, HTML) work out of the box; office/PDF formats import their
backing library lazily and raise a helpful :class:`ConversionError` when it is
not installed.
"""

from __future__ import annotations

import csv
import io
import json
import re
from html.parser import HTMLParser
from pathlib import Path

from .core import (
    ConversionError,
    ConversionResult,
    ConvertOptions,
    register,
    to_markdown_table,
)


# --------------------------------------------------------------------------- #
# Shared helpers
# --------------------------------------------------------------------------- #

def _read_text(path: Path) -> str:
    """Read a text file, tolerating common encodings."""
    data = path.read_bytes()
    for enc in ("utf-8-sig", "utf-8"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    try:  # optional, nicer guesses for legacy encodings
        import chardet  # type: ignore

        guess = chardet.detect(data)
        if guess.get("encoding"):
            return data.decode(guess["encoding"], errors="replace")
    except Exception:
        pass
    return data.decode("latin-1", errors="replace")


def _trim_matrix(rows: list[list[object]]) -> list[list[object]]:
    """Drop trailing all-empty rows and trailing all-empty columns."""

    def empty(value: object) -> bool:
        return value is None or str(value).strip() == ""

    while rows and all(empty(c) for c in rows[-1]):
        rows.pop()
    if not rows:
        return rows
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    while width > 0 and all(empty(r[width - 1]) for r in rows):
        width -= 1
        rows = [r[:width] for r in rows]
    return rows


# --------------------------------------------------------------------------- #
# Plain text / Markdown
# --------------------------------------------------------------------------- #

@register(
    ".txt", ".text", ".log", ".md", ".markdown",
    description="Plain text / Markdown",
)
def convert_text(path: Path, options: ConvertOptions) -> str:
    text = _read_text(path)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    return text.rstrip("\n") + "\n"


# --------------------------------------------------------------------------- #
# CSV / TSV
# --------------------------------------------------------------------------- #

@register(".csv", description="Comma-separated values")
def convert_csv(path: Path, options: ConvertOptions) -> str:
    return _convert_delimited(path, options, delimiter=None)


@register(".tsv", ".tab", description="Tab-separated values")
def convert_tsv(path: Path, options: ConvertOptions) -> str:
    return _convert_delimited(path, options, delimiter="\t")


def _convert_delimited(path: Path, options: ConvertOptions, delimiter: str | None) -> str:
    text = _read_text(path)
    if delimiter is None:
        try:
            delimiter = csv.Sniffer().sniff(text[:4096], delimiters=",;\t|").delimiter
        except csv.Error:
            delimiter = ","
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = [row for row in reader if row]
    if not rows:
        return "*(empty file)*\n"
    return to_markdown_table(rows, has_header=options.has_header)


# --------------------------------------------------------------------------- #
# JSON
# --------------------------------------------------------------------------- #

@register(".json", description="JSON")
def convert_json(path: Path, options: ConvertOptions) -> str:
    text = _read_text(path)
    try:
        data = json.loads(text)
    except json.JSONDecodeError as exc:
        raise ConversionError(f"Invalid JSON in {path.name}: {exc}") from exc
    return _json_to_md(data)


def _json_to_md(data: object) -> str:
    # A list of flat-ish dicts renders nicely as a table.
    if isinstance(data, list) and data and all(isinstance(x, dict) for x in data):
        columns: list[str] = []
        for item in data:
            for key in item:
                if key not in columns:
                    columns.append(key)
        rows = [columns]
        for item in data:
            rows.append([_json_scalar(item.get(col, "")) for col in columns])
        return to_markdown_table(rows, has_header=True)
    # Anything else: a pretty-printed fenced code block.
    pretty = json.dumps(data, indent=2, ensure_ascii=False)
    return f"```json\n{pretty}\n```\n"


def _json_scalar(value: object) -> str:
    if isinstance(value, (dict, list)):
        return json.dumps(value, ensure_ascii=False)
    if isinstance(value, bool):
        return "true" if value else "false"
    if value is None:
        return ""
    return str(value)


# --------------------------------------------------------------------------- #
# HTML
# --------------------------------------------------------------------------- #

@register(".html", ".htm", description="HTML")
def convert_html(path: Path, options: ConvertOptions) -> ConversionResult:
    text = _read_text(path)
    try:
        from markdownify import markdownify as _md  # type: ignore

        md = _md(text, heading_style="ATX")
        md = re.sub(r"\n{3,}", "\n\n", md).strip() + "\n"
        return ConversionResult(source=path, markdown=md)
    except ImportError:
        parser = _HtmlToMd()
        parser.feed(text)
        warning = "Used the built-in HTML converter; install 'markdownify' for higher fidelity."
        return ConversionResult(source=path, markdown=parser.get_markdown(), warnings=[warning])


class _HtmlToMd(HTMLParser):
    """A small, dependency-free HTML -> Markdown converter for common tags."""

    HEADINGS = {"h1": 1, "h2": 2, "h3": 3, "h4": 4, "h5": 5, "h6": 6}
    BLOCKS = ("p", "div", "section", "article", "blockquote", "table", "tr")

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.skip = 0                      # inside <script>/<style>/<head>
        self.pre = 0                       # inside <pre>
        self.lists: list[list] = []        # stack of ["ul"|"ol", counter]
        self.link: tuple | None = None     # (href, [text parts]) while in <a>

    # -- emit -------------------------------------------------------------- #
    def _emit(self, text: str) -> None:
        if self.link is not None:
            self.link[1].append(text)
        else:
            self.parts.append(text)

    # -- tags -------------------------------------------------------------- #
    def handle_starttag(self, tag, attrs):
        attr = dict(attrs)
        if tag in ("script", "style", "head"):
            self.skip += 1
            return
        if self.skip:
            return
        if tag in self.HEADINGS:
            self._emit("\n\n" + "#" * self.HEADINGS[tag] + " ")
        elif tag in self.BLOCKS:
            self._emit("\n\n")
        elif tag == "br":
            self._emit("  \n")
        elif tag == "hr":
            self._emit("\n\n---\n\n")
        elif tag in ("strong", "b"):
            self._emit("**")
        elif tag in ("em", "i"):
            self._emit("*")
        elif tag == "code" and not self.pre:
            self._emit("`")
        elif tag == "pre":
            self.pre += 1
            self._emit("\n\n```\n")
        elif tag in ("ul", "ol"):
            self.lists.append([tag, 0])
            self._emit("\n")
        elif tag == "li":
            indent = "  " * max(0, len(self.lists) - 1)
            if self.lists and self.lists[-1][0] == "ol":
                self.lists[-1][1] += 1
                self._emit(f"\n{indent}{self.lists[-1][1]}. ")
            else:
                self._emit(f"\n{indent}- ")
        elif tag == "a":
            self.link = (attr.get("href"), [])
        elif tag == "img":
            self._emit(f"![{attr.get('alt', '')}]({attr.get('src', '')})")

    def handle_startendtag(self, tag, attrs):
        self.handle_starttag(tag, attrs)

    def handle_endtag(self, tag):
        if tag in ("script", "style", "head"):
            self.skip = max(0, self.skip - 1)
            return
        if self.skip:
            return
        if tag in self.HEADINGS or tag in self.BLOCKS:
            self._emit("\n\n")
        elif tag in ("strong", "b"):
            self._emit("**")
        elif tag in ("em", "i"):
            self._emit("*")
        elif tag == "code" and not self.pre:
            self._emit("`")
        elif tag == "pre":
            self._emit("\n```\n\n")
            self.pre = max(0, self.pre - 1)
        elif tag in ("ul", "ol"):
            if self.lists:
                self.lists.pop()
            self._emit("\n")
        elif tag == "a" and self.link is not None:
            href, text_parts = self.link
            self.link = None
            text = "".join(text_parts).strip()
            self.parts.append(f"[{text}]({href})" if href else text)
        elif tag in ("td", "th"):
            self._emit(" ")

    def handle_data(self, data):
        if self.skip:
            return
        if self.pre:
            self._emit(data)
            return
        text = re.sub(r"\s+", " ", data)
        if not text.strip():
            if self.parts and not self.parts[-1].endswith((" ", "\n")):
                self._emit(" ")
            return
        self._emit(text)

    def get_markdown(self) -> str:
        md = "".join(self.parts)
        md = re.sub(r"[ \t]+\n", "\n", md)
        md = re.sub(r"\n{3,}", "\n\n", md)
        return md.strip() + "\n"


# --------------------------------------------------------------------------- #
# Excel (.xlsx / .xlsm)  -- requires openpyxl
# --------------------------------------------------------------------------- #

@register(".xlsx", ".xlsm", description="Excel workbook (needs openpyxl)")
def convert_xlsx(path: Path, options: ConvertOptions) -> ConversionResult:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        raise ConversionError(
            "Reading .xlsx requires 'openpyxl' — install it with: pip install openpyxl"
        )

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    blocks: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows = [
                ["" if cell is None else cell for cell in row]
                for row in sheet.iter_rows(values_only=True)
            ]
            rows = _trim_matrix(rows)
            if not rows:
                blocks.append(f"## {sheet.title}\n\n*(empty sheet)*")
                continue
            table = to_markdown_table(rows, has_header=options.has_header)
            blocks.append(f"## {sheet.title}\n\n{table.rstrip()}")
    finally:
        workbook.close()

    markdown = ("\n\n".join(blocks)).rstrip() + "\n"
    return ConversionResult(source=path, markdown=markdown)


# --------------------------------------------------------------------------- #
# PDF  -- prefers pdfplumber (text + tables), falls back to pypdf (text only)
# --------------------------------------------------------------------------- #

@register(".pdf", description="PDF document (needs pdfplumber or pypdf)")
def convert_pdf(path: Path, options: ConvertOptions) -> ConversionResult:
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        pdfplumber = None

    if pdfplumber is not None:
        return _pdf_with_plumber(path, options, pdfplumber)

    try:
        import pypdf  # type: ignore
    except ImportError:
        raise ConversionError(
            "Reading .pdf requires 'pdfplumber' (text + tables) or 'pypdf' (text only) — "
            "install with: pip install pdfplumber"
        )

    reader = pypdf.PdfReader(str(path))
    parts: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        parts.append(f"## Page {index}\n\n{text}" if options.page_breaks else text)
    markdown = ("\n\n".join(parts)).strip() + "\n"
    warning = "Used pypdf (text only); install 'pdfplumber' to also extract tables."
    return ConversionResult(source=path, markdown=markdown, warnings=[warning])


def _pdf_with_plumber(path: Path, options: ConvertOptions, pdfplumber) -> ConversionResult:
    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            section: list[str] = []
            if options.page_breaks:
                section.append(f"## Page {index}")
            text = (page.extract_text() or "").strip()
            if text:
                section.append(text)
            for table in page.extract_tables() or []:
                if table:
                    section.append(to_markdown_table(table, has_header=options.has_header).rstrip())
            if len(section) > (1 if options.page_breaks else 0):
                parts.append("\n\n".join(section))
    markdown = ("\n\n".join(parts)).strip() + "\n"
    return ConversionResult(source=path, markdown=markdown)


# --------------------------------------------------------------------------- #
# Word (.docx)  -- requires python-docx
# --------------------------------------------------------------------------- #

@register(".docx", description="Word document (needs python-docx)")
def convert_docx(path: Path, options: ConvertOptions) -> ConversionResult:
    try:
        import docx  # type: ignore
        from docx.document import Document as _Doc
        from docx.table import Table as _Table
        from docx.text.paragraph import Paragraph as _Paragraph
    except ImportError:
        raise ConversionError(
            "Reading .docx requires 'python-docx' — install it with: pip install python-docx"
        )

    document = docx.Document(str(path))
    parts: list[str] = []

    for block in _iter_docx_blocks(document, _Doc, _Paragraph, _Table):
        if isinstance(block, _Paragraph):
            rendered = _docx_paragraph(block)
            if rendered:
                parts.append(rendered)
        else:  # table
            rows = [[cell.text for cell in row.cells] for row in block.rows]
            parts.append(to_markdown_table(rows, has_header=options.has_header).rstrip())

    markdown = ("\n\n".join(parts)).strip() + "\n"
    return ConversionResult(source=path, markdown=markdown)


def _iter_docx_blocks(document, _Doc, _Paragraph, _Table):
    """Yield paragraphs and tables of a docx body in document order."""
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P

    parent = document.element.body
    for child in parent.iterchildren():
        if isinstance(child, CT_P):
            yield _Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield _Table(child, document)


def _docx_paragraph(paragraph) -> str:
    text = "".join(_docx_run(run) for run in paragraph.runs).strip()
    if not text:
        return ""
    style = (paragraph.style.name or "").lower() if paragraph.style else ""
    if style.startswith("heading"):
        match = re.search(r"(\d+)", style)
        level = int(match.group(1)) if match else 1
        return "#" * min(max(level, 1), 6) + " " + text
    if style == "title":
        return "# " + text
    if "list" in style:
        return "- " + text
    return text


def _docx_run(run) -> str:
    text = run.text
    if not text:
        return text
    if run.bold and run.italic:
        return f"***{text}***"
    if run.bold:
        return f"**{text}**"
    if run.italic:
        return f"*{text}*"
    return text


# --------------------------------------------------------------------------- #
# PowerPoint (.pptx)  -- requires python-pptx
# --------------------------------------------------------------------------- #

@register(".pptx", description="PowerPoint presentation (needs python-pptx)")
def convert_pptx(path: Path, options: ConvertOptions) -> ConversionResult:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise ConversionError(
            "Reading .pptx requires 'python-pptx' — install it with: pip install python-pptx"
        )

    presentation = Presentation(str(path))
    parts: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"## Slide {index}")
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.text.strip():
            parts.append(f"### {title_shape.text.strip()}")

        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if shape.has_table:
                table = shape.table
                rows = [[cell.text for cell in row.cells] for row in table.rows]
                parts.append(to_markdown_table(rows, has_header=options.has_header).rstrip())
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        indent = "  " * max(0, getattr(para, "level", 0) or 0)
                        parts.append(f"{indent}- {line}")

    markdown = ("\n\n".join(parts)).strip() + "\n"
    return ConversionResult(source=path, markdown=markdown)
